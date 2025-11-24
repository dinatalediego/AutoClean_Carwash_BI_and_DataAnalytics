"""
carwash_reporting.py

Capas de reporting para el proyecto Car Wash:
- PPTXReporter: genera presentación PowerPoint.
- PDFReporter: genera un informe ejecutivo en PDF (vía HTML/Markdown).
- TextReporter: genera script verbal y correo profesional.

Estas clases consumen:
- IntelligencePipeline (de carwash_pipeline)
- CarWashEconEngine (de carwash_econ)
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import List, Dict, Any, Optional
import os
import datetime as dt

import pandas as pd

# Librerías a usar en tu entorno local:
# pip install python-pptx reportlab

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas


# ==========================
# BASE REPORTER
# ==========================

@dataclass
class ReportingContext:
    """Contexto que comparten los reportes."""
    kpis: Dict[str, Dict[str, Any]]      # pipeline.kpis
    equilibrium_bullets: List[str]       # narrativa de equilibrio
    elasticity_bullets: List[str]        # narrativa de elasticidad (opcional)
    output_folder: str = "output"


# ==========================
# PPTX REPORTER
# ==========================

class PPTXReporter:
    """Genera una presentación PPTX con KPIs + narrativa."""

    def __init__(self, ctx: ReportingContext):
        self.ctx = ctx

    def build_presentation(self) -> Presentation:
        prs = Presentation()
        today = dt.date.today().strftime("%d-%b-%Y")

        # Slide 1 – Título
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "INTELIGENCIA COMERCIAL – CAR WASH"
        subtitle = slide.placeholders[1]
        subtitle.text = f"Reporte y Modelos – Actualizado al {today}"

        # Slide 2 – KPIs principales
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "KPIs principales"

        tf = slide.placeholders[1].text_frame
        ticket_kpis = self.ctx.kpis.get("ticket", {})
        repeat_kpis = self.ctx.kpis.get("repeat", {})
        forecast_kpis = self.ctx.kpis.get("forecast", {})
        pricing_kpis = self.ctx.kpis.get("pricing", {})

        tf.text = f"Ticket promedio estimado: S/ {ticket_kpis.get('ticket_promedio_estimado', 'N/D'):.2f}"
        for line in [
            f"Probabilidad media de recurrencia: {repeat_kpis.get('probabilidad_media_recurrencia', 'N/D'):.2%}"
            if isinstance(repeat_kpis.get('probabilidad_media_recurrencia', None), (float, int))
            else f"Probabilidad media de recurrencia: {repeat_kpis.get('probabilidad_media_recurrencia', 'N/D')}",
            f"Ingresos diarios esperados (forecast): S/ {forecast_kpis.get('ingresos_diarios_esperados', 'N/D')}",
            f"Elasticidad precio (placeholder): {pricing_kpis.get('elasticidad_precio', 'N/D')}"
        ]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0

        # Slide 3 – Equilibrio: “sin intervención…”
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Equilibrio de demanda (sin intervención)"
        tf = slide.placeholders[1].text_frame
        tf.text = self.ctx.equilibrium_bullets[0] if self.ctx.equilibrium_bullets else ""
        for b in self.ctx.equilibrium_bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

        # Slide 4 – Elasticidad / pricing (opcional)
        if self.ctx.elasticity_bullets:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Elasticidad precio-demanda"
            tf = slide.placeholders[1].text_frame
            tf.text = self.ctx.elasticity_bullets[0]
            for b in self.ctx.elasticity_bullets[1:]:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0

        # Slide 5 – Recomendaciones de corto plazo
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Recomendaciones de corto plazo"
        tf = slide.placeholders[1].text_frame
        tf.text = "1) Activar upsell Standar → Premium en segmentos de baja elasticidad (camionetas)."
        for line in [
            "2) Promociones tácticas en días valle para elevar demanda por debajo del equilibrio.",
            "3) Pruebas controladas de precio en Standar para medir elasticidad real por segmento.",
            "4) Programa simple de fidelización para clientes recurrentes (placas con visitas ≥2)."
        ]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0

        return prs

    def save(self, filename: str) -> str:
        os.makedirs(self.ctx.output_folder, exist_ok=True)
        full_path = os.path.join(self.ctx.output_folder, filename)
        prs = self.build_presentation()
        prs.save(full_path)
        return full_path


# ==========================
# PDF REPORTER
# ==========================

class PDFReporter:
    """
    Informe ejecutivo breve en PDF.
    Versión simple: texto plano en una página usando reportlab.
    """

    def __init__(self, ctx: ReportingContext):
        self.ctx = ctx

    def save(self, filename: str) -> str:
        os.makedirs(self.ctx.output_folder, exist_ok=True)
        full_path = os.path.join(self.ctx.output_folder, filename)

        c = canvas.Canvas(full_path, pagesize=A4)
        width, height = A4

        c.setFont("Helvetica-Bold", 16)
        c.drawString(50, height - 50, "Informe Ejecutivo – Inteligencia Comercial Car Wash")

        c.setFont("Helvetica", 11)
        y = height - 90

        # secciones clave
        sections = [
            ("Equilibrio de demanda", self.ctx.equilibrium_bullets),
            ("Elasticidad y pricing", self.ctx.elasticity_bullets),
        ]

        for title, bullets in sections:
            if not bullets:
                continue
            c.setFont("Helvetica-Bold", 12)
            c.drawString(50, y, title)
            y -= 18
            c.setFont("Helvetica", 11)
            for b in bullets:
                for line in self._wrap_text(b, max_chars=90):
                    c.drawString(60, y, f"- {line}")
                    y -= 14
                    if y < 80:
                        c.showPage()
                        y = height - 80
                        c.setFont("Helvetica", 11)
            y -= 10

        c.showPage()
        c.save()
        return full_path

    @staticmethod
    def _wrap_text(text: str, max_chars: int = 90) -> List[str]:
        """Divide una línea larga en varias para caber en la página."""
        words = text.split()
        lines = []
        current = []
        current_len = 0
        for w in words:
            if current_len + len(w) + 1 > max_chars:
                lines.append(" ".join(current))
                current = [w]
                current_len = len(w)
            else:
                current.append(w)
                current_len += len(w) + 1
        if current:
            lines.append(" ".join(current))
        return lines


# ==========================
# TEXT REPORTER
# ==========================

class TextReporter:
    """
    Genera:
    - script verbal (para masterclass / exposición)
    - borrador de correo profesional para propuesta de consultoría
    """

    def __init__(self, ctx: ReportingContext):
        self.ctx = ctx

    def build_script_verbal(self) -> str:
        partes = []
        partes.append("Script verbal – Cierre de presentación (Car Wash, microeconomía + econometría)\n")

        # Equilibrio
        partes.append("1) Equilibrio de demanda (sin intervención):")
        for b in self.ctx.equilibrium_bullets:
            partes.append(f"   - {b}")

        # Elasticidad
        if self.ctx.elasticity_bullets:
            partes.append("\n2) Elasticidad precio-demanda:")
            for b in self.ctx.elasticity_bullets:
                partes.append(f"   - {b}")

        partes.append("\n3) Cierre propositivo:")
        partes.append("   - \"Sin hacer nada, el negocio se queda estable en este equilibrio;\"")
        partes.append("   - \"con intervención inteligente, podemos elevar los valles, proteger los picos y mover el ticket promedio.\"")
        partes.append("   - \"Mi propuesta es empezar con un sprint de corto plazo para demostrar este impacto y, si lo ven, construir juntos el modelo de largo plazo.\"")

        return "\n".join(partes)

    def build_correo_propuesta(self, nombre_cliente: str = "Equipo") -> str:
        correo = f"""
Estimado {nombre_cliente},

A partir del análisis realizado, hoy sabemos que:

- La demanda del Car Wash, sin intervención, se estabiliza alrededor de un nivel de equilibrio diario.
- Existen picos muy claros de alta utilización y valles donde el potencial del negocio no se está capturando.
- La elasticidad precio-demanda y el comportamiento de los clientes recurrentes nos dan espacio real para mejorar ticket promedio, recurrencia y utilización de la capacidad.

Mi propuesta es trabajar en dos etapas:

1) Un sprint de corto plazo (4–5 semanas) para:
   - Afinar el modelo de forecast y elasticidad,
   - Diseñar e implementar una primera batería de intervenciones en precios y promociones,
   - Medir el impacto en ticket y en llenado de días valle.

2) Si el impacto es claro y medible, pasar a un esquema mensual de Inteligencia Comercial continua,
   donde tengamos modelado en forma permanente el equilibrio, el pricing y la recurrencia de clientes.

Quedo atento para coordinar una reunión breve y revisar juntos los resultados y próximos pasos.

Saludos,

[Tu nombre]
[Tu rol]
"""
        return correo.strip()

    def save_texts(self, script_filename: str, correo_filename: str) -> Dict[str, str]:
        os.makedirs(self.ctx.output_folder, exist_ok=True)
        paths = {}

        script_path = os.path.join(self.ctx.output_folder, script_filename)
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(self.build_script_verbal())
        paths["script"] = script_path

        correo_path = os.path.join(self.ctx.output_folder, correo_filename)
        with open(correo_path, "w", encoding="utf-8") as f:
            f.write(self.build_correo_propuesta())
        paths["correo"] = correo_path

        return paths
